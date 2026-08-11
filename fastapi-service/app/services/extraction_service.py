import csv
import io
import json
import subprocess
import tempfile
from html.parser import HTMLParser
from pathlib import Path
from xml.etree import ElementTree


class _HtmlTextExtractor(HTMLParser):
    skip_tags = {"script", "style", "nav", "footer", "header"}

    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []
        self.skip_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag.lower() in self.skip_tags:
            self.skip_depth += 1

    def handle_endtag(self, tag: str) -> None:
        if tag.lower() in self.skip_tags and self.skip_depth > 0:
            self.skip_depth -= 1

    def handle_data(self, data: str) -> None:
        value = data.strip()

        if self.skip_depth == 0 and value:
            self.parts.append(value)


class ExtractionService:
    max_text_chars = 500_000

    def extract_text(self, file_bytes: bytes, extension: str) -> str:
        normalized = extension.lower().lstrip(".")
        extractors = {
            "pdf": self._extract_pdf,
            "docx": self._extract_docx,
            "doc": self._extract_doc,
            "pptx": self._extract_pptx,
            "ppt": self._extract_ppt,
            "xlsx": self._extract_xlsx,
            "csv": self._extract_csv,
            "txt": self._extract_txt,
            "html": self._extract_html,
            "xml": self._extract_xml,
            "json": self._extract_json,
            "png": self._extract_image,
            "jpg": self._extract_image,
            "jpeg": self._extract_image,
        }
        extractor = extractors.get(normalized)

        if extractor is None:
            raise ValueError(f"Unsupported file type for RAG: .{extension}")

        text = extractor(file_bytes, normalized)
        return text[: self.max_text_chars]

    def _extract_pdf(self, file_bytes: bytes, _: str) -> str:
        import pymupdf4llm

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=True) as temp_file:
            temp_file.write(file_bytes)
            temp_file.flush()
            return pymupdf4llm.to_markdown(temp_file.name)

    def _extract_docx(self, file_bytes: bytes, _: str = "docx") -> str:
        from docx import Document

        document = Document(io.BytesIO(file_bytes))
        parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]

        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))

        return "\n".join(parts)

    def _extract_doc(self, file_bytes: bytes, extension: str) -> str:
        return self._extract_docx(self._convert_legacy_office(file_bytes, extension, "docx"))

    def _extract_pptx(self, file_bytes: bytes, _: str = "pptx") -> str:
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(file_bytes))
        slides: list[str] = []

        for index, slide in enumerate(presentation.slides, start=1):
            parts = [f"--- Slide {index} ---"]

            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text_frame.text.strip()
                    if text:
                        parts.append(text)

            slides.append("\n".join(parts))

        return "\n\n".join(slides)

    def _extract_ppt(self, file_bytes: bytes, extension: str) -> str:
        return self._extract_pptx(self._convert_legacy_office(file_bytes, extension, "pptx"))

    def _extract_xlsx(self, file_bytes: bytes, _: str) -> str:
        from openpyxl import load_workbook

        workbook = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        output: list[str] = []

        for sheet in workbook.worksheets:
            output.append(f"--- Sheet: {sheet.title} ---")
            headers: list[str] = []

            for row_index, row in enumerate(sheet.iter_rows(values_only=True)):
                values = ["" if value is None else str(value) for value in row]

                if row_index == 0:
                    headers = values
                    output.append(" | ".join(value for value in values if value))
                    continue

                row_parts = []
                for cell_index, value in enumerate(values):
                    if not value:
                        continue

                    label = headers[cell_index] if cell_index < len(headers) else f"Column {cell_index + 1}"
                    row_parts.append(f"{label}={value}" if label else value)

                if row_parts:
                    output.append(", ".join(row_parts))

        return "\n".join(output)

    def _extract_csv(self, file_bytes: bytes, _: str) -> str:
        text = file_bytes.decode("utf-8", errors="replace")
        reader = csv.DictReader(io.StringIO(text))
        rows = []

        for row in reader:
            rows.append(", ".join(f"{key}={value}" for key, value in row.items() if value))

        return "\n".join(rows) if rows else text

    def _extract_txt(self, file_bytes: bytes, _: str) -> str:
        return file_bytes.decode("utf-8", errors="replace")

    def _extract_html(self, file_bytes: bytes, _: str) -> str:
        parser = _HtmlTextExtractor()
        parser.feed(file_bytes.decode("utf-8", errors="replace"))
        return "\n".join(parser.parts)

    def _extract_xml(self, file_bytes: bytes, _: str) -> str:
        root = ElementTree.fromstring(file_bytes.decode("utf-8", errors="replace"))
        return " ".join(part.strip() for part in root.itertext() if part.strip())

    def _extract_json(self, file_bytes: bytes, _: str) -> str:
        data = json.loads(file_bytes.decode("utf-8", errors="replace"))
        flattened: list[str] = []

        def visit(value: object, prefix: str = "") -> None:
            if isinstance(value, dict):
                for key, child in value.items():
                    visit(child, f"{prefix}.{key}" if prefix else str(key))
            elif isinstance(value, list):
                for index, child in enumerate(value):
                    visit(child, f"{prefix}[{index}]")
            else:
                flattened.append(f"{prefix}: {value}" if prefix else str(value))

        visit(data)
        return "\n".join(flattened)

    def _extract_image(self, file_bytes: bytes, extension: str) -> str:
        import fitz

        document = fitz.open(stream=file_bytes, filetype=extension)
        page = document[0]
        text = page.get_text("text")

        if not text.strip():
            text_page = page.get_textpage_ocr(language="eng")
            text = page.get_text("text", textpage=text_page)

        return text

    def _convert_legacy_office(
        self,
        file_bytes: bytes,
        source_extension: str,
        target_extension: str,
    ) -> bytes:
        with tempfile.TemporaryDirectory() as temp_dir:
            source_path = Path(temp_dir) / f"input.{source_extension}"
            source_path.write_bytes(file_bytes)

            subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    target_extension,
                    "--outdir",
                    temp_dir,
                    str(source_path),
                ],
                check=True,
                timeout=45,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
            )

            converted_path = Path(temp_dir) / f"input.{target_extension}"
            if not converted_path.exists():
                raise ValueError("LibreOffice did not produce a converted preview file.")

            return converted_path.read_bytes()
